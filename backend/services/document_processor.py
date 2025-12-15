"""
Document Processor - Pipeline completo para RAG
Processa PPTX, PDF, DOCX, TXT e extrai texto estruturado
"""

import os
import re
import uuid
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass
from datetime import datetime
import traceback


# ============================================================
# 📄 ESTRUTURAS DE DADOS
# ============================================================

@dataclass
class DocumentChunk:
    """Representa um chunk de documento para indexação"""
    chunk_id: str
    text: str
    metadata: Dict
    source_file: str
    chunk_index: int
    total_chunks: int


@dataclass  
class ProcessingResult:
    """Resultado do processamento de um documento"""
    success: bool
    document_id: str
    filename: str
    file_type: str
    chunks_created: int
    error_message: Optional[str] = None
    processing_time_ms: int = 0


# ============================================================
# 📄 EXTRATORES DE TEXTO POR TIPO
# ============================================================

def extract_text_from_pptx(file_path: str) -> List[Dict]:
    """
    Extrai texto de arquivo PowerPoint (.pptx)
    Retorna lista de dicionários com texto e metadata por slide
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        slides_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
            
            if slide_text_parts:
                slides_content.append({
                    "text": "\n".join(slide_text_parts),
                    "metadata": {
                        "slide_number": slide_num,
                        "content_type": "slide"
                    }
                })
        
        return slides_content
        
    except ImportError:
        raise ImportError("python-pptx não instalado. Execute: pip install python-pptx")
    except Exception as e:
        raise Exception(f"Erro ao processar PPTX: {str(e)}")


def extract_text_from_pdf(file_path: str) -> List[Dict]:
    """
    Extrai texto de arquivo PDF
    Retorna lista de dicionários com texto e metadata por página
    """
    try:
        from PyPDF2 import PdfReader
        
        reader = PdfReader(file_path)
        pages_content = []
        
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                pages_content.append({
                    "text": text.strip(),
                    "metadata": {
                        "page_number": page_num,
                        "content_type": "page"
                    }
                })
        
        return pages_content
        
    except ImportError:
        raise ImportError("PyPDF2 não instalado. Execute: pip install PyPDF2")
    except Exception as e:
        raise Exception(f"Erro ao processar PDF: {str(e)}")


def extract_text_from_docx(file_path: str) -> List[Dict]:
    """
    Extrai texto de arquivo Word (.docx)
    Retorna lista de dicionários com texto e metadata por seção/parágrafo
    """
    try:
        from docx import Document
        
        doc = Document(file_path)
        content = []
        
        # Extrai parágrafos
        current_section = []
        section_num = 1
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
                
            # Detecta títulos/headers para criar seções
            if para.style and para.style.name.startswith('Heading'):
                # Salva seção anterior se existir
                if current_section:
                    content.append({
                        "text": "\n".join(current_section),
                        "metadata": {
                            "section_number": section_num,
                            "content_type": "section"
                        }
                    })
                    section_num += 1
                    current_section = []
                
                current_section.append(text)
            else:
                current_section.append(text)
        
        # Adiciona última seção
        if current_section:
            content.append({
                "text": "\n".join(current_section),
                "metadata": {
                    "section_number": section_num,
                    "content_type": "section"
                }
            })
        
        # Se não encontrou estrutura, retorna tudo como um bloco
        if not content:
            full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            if full_text:
                content.append({
                    "text": full_text,
                    "metadata": {
                        "section_number": 1,
                        "content_type": "full_document"
                    }
                })
        
        return content
        
    except ImportError:
        raise ImportError("python-docx não instalado. Execute: pip install python-docx")
    except Exception as e:
        raise Exception(f"Erro ao processar DOCX: {str(e)}")


def extract_text_from_txt(file_path: str) -> List[Dict]:
    """
    Extrai texto de arquivo TXT
    Retorna lista com o conteúdo do arquivo
    """
    try:
        # Tenta diferentes encodings
        encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                    if text.strip():
                        return [{
                            "text": text.strip(),
                            "metadata": {
                                "encoding": encoding,
                                "content_type": "text_file"
                            }
                        }]
                    return []
            except UnicodeDecodeError:
                continue
        
        raise Exception("Não foi possível decodificar o arquivo com nenhum encoding conhecido")
        
    except Exception as e:
        raise Exception(f"Erro ao processar TXT: {str(e)}")


# ============================================================
# 📄 NORMALIZAÇÃO DE TEXTO
# ============================================================

def normalize_text(text: str) -> str:
    """
    Normaliza texto para melhor indexação:
    - Remove caracteres de controle
    - Normaliza espaços em branco
    - Remove linhas vazias excessivas
    """
    if not text:
        return ""
    
    # Remove caracteres de controle (exceto newlines e tabs)
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text)
    
    # Substitui múltiplos espaços por um único
    text = re.sub(r'[ \t]+', ' ', text)
    
    # Substitui múltiplas linhas vazias por no máximo duas
    text = re.sub(r'\n{3,}', '\n\n', text)
    
    # Remove espaços no início e fim de cada linha
    lines = [line.strip() for line in text.split('\n')]
    text = '\n'.join(lines)
    
    return text.strip()


# ============================================================
# 📄 CHUNKING DE TEXTO
# ============================================================

def chunk_text(
    text: str, 
    chunk_size: int = 500, 
    overlap: int = 50,
    min_chunk_size: int = 100
) -> List[str]:
    """
    Divide texto em chunks menores para melhor indexação.
    
    Args:
        text: Texto a ser dividido
        chunk_size: Tamanho aproximado de cada chunk em caracteres
        overlap: Quantidade de caracteres de sobreposição entre chunks
        min_chunk_size: Tamanho mínimo para criar um chunk
        
    Returns:
        Lista de chunks de texto
    """
    if not text or len(text) < min_chunk_size:
        return [text] if text else []
    
    # Se texto é menor que chunk_size, retorna inteiro
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        # Define o fim do chunk
        end = start + chunk_size
        
        # Se não é o último chunk, tenta cortar em um ponto natural
        if end < len(text):
            # Procura por quebra de parágrafo primeiro
            paragraph_break = text.rfind('\n\n', start, end)
            if paragraph_break > start + min_chunk_size:
                end = paragraph_break
            else:
                # Procura por quebra de linha
                line_break = text.rfind('\n', start, end)
                if line_break > start + min_chunk_size:
                    end = line_break
                else:
                    # Procura por ponto final
                    sentence_end = text.rfind('. ', start, end)
                    if sentence_end > start + min_chunk_size:
                        end = sentence_end + 1
        
        chunk = text[start:end].strip()
        if len(chunk) >= min_chunk_size:
            chunks.append(chunk)
        
        # Próximo chunk começa com overlap
        start = end - overlap if end < len(text) else len(text)
    
    return chunks


# ============================================================
# 📄 PROCESSADOR PRINCIPAL
# ============================================================

SUPPORTED_EXTENSIONS = {
    '.pptx': extract_text_from_pptx,
    '.pdf': extract_text_from_pdf,
    '.docx': extract_text_from_docx,
    '.txt': extract_text_from_txt,
}


def get_file_extension(filename: str) -> str:
    """Retorna extensão do arquivo em lowercase"""
    return os.path.splitext(filename)[1].lower()


def is_supported_file(filename: str) -> bool:
    """Verifica se o tipo de arquivo é suportado"""
    ext = get_file_extension(filename)
    return ext in SUPPORTED_EXTENSIONS


def process_document(
    file_path: str,
    filename: str,
    chunk_size: int = 500,
    chunk_overlap: int = 50
) -> Tuple[List[DocumentChunk], ProcessingResult]:
    """
    Processa um documento e retorna chunks prontos para indexação.
    
    Pipeline completo:
    1. Identifica tipo de arquivo
    2. Extrai texto bruto
    3. Normaliza texto
    4. Divide em chunks
    5. Adiciona metadata
    
    Args:
        file_path: Caminho do arquivo
        filename: Nome original do arquivo
        chunk_size: Tamanho dos chunks
        chunk_overlap: Sobreposição entre chunks
        
    Returns:
        Tupla com (lista de chunks, resultado do processamento)
    """
    start_time = datetime.now()
    document_id = str(uuid.uuid4())
    
    # Valida extensão
    ext = get_file_extension(filename)
    if ext not in SUPPORTED_EXTENSIONS:
        return [], ProcessingResult(
            success=False,
            document_id=document_id,
            filename=filename,
            file_type=ext,
            chunks_created=0,
            error_message=f"Tipo de arquivo não suportado: {ext}. Suportados: {list(SUPPORTED_EXTENSIONS.keys())}"
        )
    
    try:
        # 1. Extrai texto
        extractor = SUPPORTED_EXTENSIONS[ext]
        raw_content = extractor(file_path)
        
        if not raw_content:
            return [], ProcessingResult(
                success=False,
                document_id=document_id,
                filename=filename,
                file_type=ext,
                chunks_created=0,
                error_message="Nenhum conteúdo de texto encontrado no arquivo"
            )
        
        # 2. Processa cada bloco de conteúdo
        all_chunks = []
        chunk_index = 0
        
        for content_block in raw_content:
            # Normaliza texto
            normalized_text = normalize_text(content_block["text"])
            
            if not normalized_text:
                continue
            
            # Divide em chunks
            text_chunks = chunk_text(
                normalized_text, 
                chunk_size=chunk_size, 
                overlap=chunk_overlap
            )
            
            # Cria DocumentChunk para cada pedaço
            for chunk_text_content in text_chunks:
                chunk = DocumentChunk(
                    chunk_id=f"{document_id}_{chunk_index}",
                    text=chunk_text_content,
                    metadata={
                        **content_block.get("metadata", {}),
                        "document_id": document_id,
                        "filename": filename,
                        "file_type": ext,
                        "processed_at": datetime.utcnow().isoformat()
                    },
                    source_file=filename,
                    chunk_index=chunk_index,
                    total_chunks=0  # Será atualizado depois
                )
                all_chunks.append(chunk)
                chunk_index += 1
        
        # Atualiza total_chunks em todos
        for chunk in all_chunks:
            chunk.total_chunks = len(all_chunks)
        
        # Calcula tempo de processamento
        processing_time = int((datetime.now() - start_time).total_seconds() * 1000)
        
        return all_chunks, ProcessingResult(
            success=True,
            document_id=document_id,
            filename=filename,
            file_type=ext,
            chunks_created=len(all_chunks),
            processing_time_ms=processing_time
        )
        
    except Exception as e:
        processing_time = int((datetime.now() - start_time).total_seconds() * 1000)
        return [], ProcessingResult(
            success=False,
            document_id=document_id,
            filename=filename,
            file_type=ext,
            chunks_created=0,
            error_message=f"{str(e)}\n{traceback.format_exc()}",
            processing_time_ms=processing_time
        )


# ============================================================
# 📄 FUNÇÕES DE CONVENIÊNCIA
# ============================================================

def get_supported_extensions() -> List[str]:
    """Retorna lista de extensões suportadas"""
    return list(SUPPORTED_EXTENSIONS.keys())


def validate_file_for_processing(filename: str, file_size_bytes: int, max_size_mb: int = 50) -> Tuple[bool, str]:
    """
    Valida se arquivo pode ser processado.
    
    Args:
        filename: Nome do arquivo
        file_size_bytes: Tamanho em bytes
        max_size_mb: Tamanho máximo em MB
        
    Returns:
        Tupla (é válido, mensagem de erro se não válido)
    """
    # Valida extensão
    ext = get_file_extension(filename)
    if ext not in SUPPORTED_EXTENSIONS:
        return False, f"Tipo não suportado: {ext}. Use: {', '.join(SUPPORTED_EXTENSIONS.keys())}"
    
    # Valida tamanho
    max_size_bytes = max_size_mb * 1024 * 1024
    if file_size_bytes > max_size_bytes:
        return False, f"Arquivo muito grande: {file_size_bytes / 1024 / 1024:.1f}MB. Máximo: {max_size_mb}MB"
    
    return True, ""

